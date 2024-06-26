__all__ = [
    'PPTXVisual',
]

import io
from typing import Any

import cv2
import numpy as np
import numpy.typing as npt
import pptx
import pptx.dml.color
import pptx.enum.shapes
import pptx.enum.text
import pptx.parts.image
import pptx.presentation
import pptx.shapes.autoshape
import pptx.shapes.picture
import pptx.shapes.shapetree
import pptx.slide
import pptx.util

from ..bases.configs import Config
from ..colors import RGB, Color
from ..registries import VisualRegistry
from .base import BaseVisual


@VisualRegistry.register_()
class PPTXVisual(BaseVisual):
    """Visualize data in the format of PowerPoint.

    The PowerPoint contains only one slide.
    For more details, refer to python-pptx_.

    .. _python-pptx: https://github.com/scanny/python-pptx
    """

    def __init__(self, width: int, height: int, **kwargs) -> None:
        """Initialize the PowerPoint with a single slide.

        To initialize a PowerPoint with width 640pt and height 426pt, use the
        following code:

            >>> visual = PPTXVisual(640, 426)

        Once initialized, the width and height of the slide cannot be altered.
        We can read the width and height of the PowerPoint by:

            >>> visual.width
            640
            >>> visual.height
            426

        Args:
            width: the width of the PowerPoint in point
            height: the height of the PowerPoint in point
        """
        self._presentation = pptx.Presentation(**kwargs)

        self.presentation.slide_width = pptx.util.Pt(width)
        self.presentation.slide_height = pptx.util.Pt(height)

        slides: pptx.slide.Slides = self.presentation.slides
        slides.add_slide(self.presentation.slide_layouts[6])

    @property
    def width(self) -> int:
        """Width of the PowerPoint."""
        width: pptx.util.Pt = self.presentation.slide_width
        return int(width.pt)

    @property
    def height(self) -> int:
        """Height of the PowerPoint."""
        height: pptx.util.Pt = self.presentation.slide_height
        return int(height.pt)

    @property
    def presentation(self) -> pptx.presentation.Presentation:
        return self._presentation

    @property
    def slide(self) -> pptx.slide.Slide:
        return self.presentation.slides[0]

    @property
    def shapes(self) -> pptx.shapes.shapetree.SlideShapes:
        return self.slide.shapes

    @classmethod
    def _set_color_format_rgb(
        cls,
        cf: pptx.dml.color.ColorFormat,
        color: Color,
    ) -> None:
        cf.rgb = pptx.dml.color.RGBColor(
            *color.to(RGB).to_tuple(),
        )

    def save(self, path: Any) -> None:
        """Save the PowerPoint.

        The save target can either be a filepath, for example:

            >>> import tempfile
            >>> with tempfile.NamedTemporaryFile() as f:
            ...     PPTXVisual(640, 426).save(f.name)

        Or it can simply be a file-like object:

            >>> with tempfile.TemporaryFile() as f:
            ...     PPTXVisual(640, 426).save(f)

        Args:
            path: destination path
        """
        self.presentation.save(path)

    def image(
        self,
        image: npt.NDArray[np.uint8],
        left: int = 0,
        top: int = 0,
        width: int | None = None,
        height: int | None = None,
        opacity: float = 1.0,
    ) -> pptx.shapes.picture.Picture:
        """Add an image to the PowerPoint.

        Suppose the image is :math:`(426, 640)`:

            >>> image = np.random.randint(0, 256, (426, 640, 3))

        In most cases, the PowerPoint is of the same size as the image, so
        that the image covers the whole background:

            >>> h, w, _ = image.shape
            >>> visual = PPTXVisual(w, h)
            >>> visual.image(image)
            <pptx.shapes.picture.Picture object at ...>

        The returned `pptx.shapes.picture.Picture` object can be used to
        fine-tune the properties of the image.
        Note that the DPI of the image should not be changed thoughtlessly,
        because of the bizarre measurements in PowerPoint.
        The most common measurement unit in PowerPoint is *Point* (pt), where
        1pt equals 1/72 inches.
        However, images are measured in pixels and 1 pixel equals 1/DPI inches.
        By default, DPI is 72 so that 1 pixel is 1pt.
        When DPI is set to other values, the size of the image may become
        larger or smaller than it is supposed to be.

        Args:
            image: :math:`(H, W, 3)`
            left: x coordinate of the left side of the image
            top: y coordinate of the top side of the image
            width: width of the image
            height: height of the image
            opacity: opacity of the image
        """
        assert 0.0 <= opacity <= 1.0

        h, w, c = image.shape
        assert c == 3

        alpha = np.ones((h, w, 1)) * 255 * opacity
        image = np.concatenate([image, alpha], axis=-1)
        success, image = cv2.imencode('.png', image)
        assert success

        picture = self.shapes.add_picture(
            io.BytesIO(image.tobytes()),
            pptx.util.Pt(left),
            pptx.util.Pt(top),
            width if width is None else pptx.util.Pt(width),
            height if height is None else pptx.util.Pt(height),
        )

        picture_image: pptx.parts.image.Image = picture.image
        assert picture_image.dpi == (72, 72)

        return picture

    def rectangle(
        self,
        left: int,
        top: int,
        width: int,
        height: int,
        color: Color = RGB(0., 0., 0.),  # noqa: B008
        thickness: int = 1,
        fill: Color | None = None,
    ) -> pptx.shapes.autoshape.Shape:
        rectangle: pptx.shapes.autoshape.Shape = self.shapes.add_shape(
            pptx.enum.shapes.MSO_AUTO_SHAPE_TYPE.RECTANGLE,  # noqa: E501 pylint: disable=no-member
            pptx.util.Pt(left),
            pptx.util.Pt(top),
            pptx.util.Pt(width),
            pptx.util.Pt(height),
        )

        line: pptx.shapes.autoshape.LineFormat = rectangle.line
        line.width = pptx.util.Pt(thickness)
        self._set_color_format_rgb(line.color, color)

        rectangle_fill: pptx.shapes.autoshape.FillFormat = rectangle.fill
        if fill is None:
            rectangle_fill.background()
        else:
            rectangle_fill.solid()
            self._set_color_format_rgb(rectangle_fill.fore_color, fill)

        return rectangle

    def text(
        self,
        text: str,
        x: int,
        y: int,
        color: Color = RGB(0., 0., 0.),  # noqa: B008
        font: Config | None = None,
    ) -> pptx.shapes.autoshape.Shape:
        if font is None:
            font = Config()

        textbox: pptx.shapes.autoshape.Shape = self.shapes.add_textbox(
            pptx.util.Pt(x),
            pptx.util.Pt(y),
            1,
            1,
        )

        text_frame = textbox.text_frame
        text_frame.text = text
        text_frame.margin_left = 0
        text_frame.margin_top = 0
        text_frame.margin_right = 0
        text_frame.margin_bottom = 0

        paragraph_font = text_frame.paragraphs[0].font
        paragraph_font.name = font.get('name', 'Times New Roman')
        paragraph_font.size = pptx.util.Pt(font.get('size', 12))
        self._set_color_format_rgb(paragraph_font.color, color)

        return textbox

    def point(
        self,
        x: int,
        y: int,
        size: int,
        color: Color = RGB(0., 0., 0.),  # noqa: B008
    ) -> Any:
        raise NotImplementedError

    def marker(
        self,
        x: int,
        y: int,
        size: int,
        color: Color = RGB(0., 0., 0.),  # noqa: B008
    ) -> Any:
        raise NotImplementedError

    def line(
        self,
        x1: int,
        y1: int,
        x2: int,
        y2: int,
        color: Color = RGB(0., 0., 0.),  # noqa: B008
        thickness: int = 1,
    ) -> Any:
        raise NotImplementedError
